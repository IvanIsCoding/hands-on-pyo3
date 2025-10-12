# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "pdf2image",
#     "python-pptx",
#     "pillow",
# ]
# ///

import os

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

# Convert PDF to images
images = convert_from_path('slides.pdf', dpi=300)

# Create PowerPoint
prs = Presentation()
blank_layout = prs.slide_layouts[6]  # blank layout

for img in images:
    slide = prs.slides.add_slide(blank_layout)
    img_path = f'temp_slide_{images.index(img)}.png'
    img.save(img_path, 'PNG')
    
    # Add image to fill slide (adjust dimensions as needed)
    slide.shapes.add_picture(img_path, 0, 0, 
                            width=prs.slide_width, 
                            height=prs.slide_height)

prs.save('slides.pptx')

for img in images:
    img_path = f'temp_slide_{images.index(img)}.png'
    try:
        os.remove(img_path)
    except OSError:
        pass